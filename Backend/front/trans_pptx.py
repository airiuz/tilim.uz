from front.translit_text import to_cyrillic, to_latin
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import re
import shutil
from zipfile import ZipFile, ZIP_DEFLATED


def transliterate(file_path1, file_path2, t):
    mytransliterate = lambda text: to_cyrillic(text) if t == '1' else to_latin(text)

    def convert(match_obj):
        data = match_obj.groupdict()
        return b"<a:t " + data['attrs'] + b">" + mytransliterate(data['text'].decode("utf-8")).encode() + b"</a:t>"

    with ZipFile(file_path1, 'r') as inzip, ZipFile(file_path2, 'w', ZIP_DEFLATED) as outzip:
        i_slides = [x.filename for x in inzip.infolist() if re.match(r'ppt\/slides\/slide\d+\.xml', x.filename)]
        for inzipinfo in inzip.infolist():
            with inzip.open(inzipinfo) as infile:
                if inzipinfo.filename in i_slides:
                    # print(re.match(rb"<a:t(?P<attrs>[^>]*)>(?P<text>[^<]+)<\/a:t>", infile.read()))
                    new_content = re.sub(rb"<a:t(?P<attrs>[^>]*)>(?P<text>[^<]+)<\/a:t>", convert, infile.read())
                    outzip.writestr(inzipinfo.filename, new_content)
                else:
                    outzip.writestr(inzipinfo.filename, infile.read())

# def transliterate(file_path, file_path2, t):

# 	def replace_text(shapes):
# 		mytransliterate = lambda text: to_cyrillic(text) if t == '1' else to_latin(text) 
# 		for shape in shapes:
# 			if shape.has_text_frame:
# 				text_frame = shape.text_frame
# 				for paragraph in text_frame.paragraphs:
# 					new_text = mytransliterate(paragraph.text)
# 					paragraph.text = new_text

# 			elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
# 				for sh in shape.shapes:
# 					if sh.has_text_frame:
# 						text_frame = sh.text_frame
# 						for paragraph in text_frame.paragraphs:
# 							new_text = mytransliterate(paragraph.text)
# 							paragraph.text = new_text

# 					if sh.has_table:
# 						for row in shape.table.rows:
# 							for cell in row.cells:
# 								cur_text = cell.text
# 								new_text = mytransliterate(cur_text)
# 								cell.text = new_text

# 			if shape.has_table:
# 				for row in shape.table.rows:
# 					for cell in row.cells:
# 						cur_text = cell.text
# 						new_text = mytransliterate(cur_text)
# 						cell.text = new_text

# 	prs = Presentation(file_path)

# 	slides = [slide for slide in prs.slides]
# 	shapes = []
# 	for slide in slides:
# 		for shape in slide.shapes:
# 			shapes.append(shape)

# 	replace_text(shapes)
# 	prs.save(file_path2)
